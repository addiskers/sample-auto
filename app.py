from flask import Flask, render_template, request, jsonify, send_file
import os
import json
import ast
import re
import traceback
import logging
import logging.handlers
import time
from datetime import datetime
from werkzeug.utils import secure_filename
from dotenv import load_dotenv
from openai import OpenAI
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_LEGEND_POSITION,
    XL_TICK_LABEL_POSITION,
    XL_TICK_MARK,
)
from enum import Enum
from typing import Optional, Dict, List, Any
import concurrent.futures
import threading

def setup_logging():
    """Setup logging configuration"""
    log_dir = '/app/logs' if os.path.exists('/app/logs') else 'logs'
    os.makedirs(log_dir, exist_ok=True)
    
    app_log_file = os.path.join(log_dir, 'app.log')
    timing_log_file = os.path.join(log_dir, 'timing.log')
    
    logging.basicConfig(level=logging.INFO)
    
    main_logger = logging.getLogger('ppt_generator')
    main_logger.setLevel(logging.INFO)
    
    timing_logger = logging.getLogger('timing')
    timing_logger.setLevel(logging.INFO)
    
    if not main_logger.handlers:
        formatter = logging.Formatter(
            '%(asctime)s - %(name)s - %(levelname)s - %(message)s',
            datefmt='%Y-%m-%d %H:%M:%S'
        )
        
        app_handler = logging.handlers.RotatingFileHandler(
            app_log_file, maxBytes=10*1024*1024, backupCount=5
        )
        app_handler.setFormatter(formatter)
        main_logger.addHandler(app_handler)
        
        console_handler = logging.StreamHandler()
        console_handler.setFormatter(formatter)
        main_logger.addHandler(console_handler)
    
    if not timing_logger.handlers:
        timing_formatter = logging.Formatter(
            '%(asctime)s - %(message)s',
            datefmt='%Y-%m-%d %H:%M:%S'
        )
        
        timing_handler = logging.handlers.RotatingFileHandler(
            timing_log_file, maxBytes=5*1024*1024, backupCount=3
        )
        timing_handler.setFormatter(timing_formatter)
        timing_logger.addHandler(timing_handler)
    
    return main_logger, timing_logger

logger, timing_logger = setup_logging()

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 20 * 1024 * 1024  
app.config['UPLOAD_FOLDER'] = 'generated_ppts'
app.config['TEMPLATE_FOLDER'] = 'templates'

os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['TEMPLATE_FOLDER'], exist_ok=True)

logger.info("Application starting up")
logger.info(f"Upload folder: {app.config['UPLOAD_FOLDER']}")
logger.info(f"Template folder: {app.config['TEMPLATE_FOLDER']}")

load_dotenv()
api_key_openAI = os.getenv("OPENAI_API_KEY")
api_key_gemini = os.getenv("GEMINI_API_KEY")

if api_key_openAI:
    logger.info("OpenAI API key loaded successfully")
else:
    logger.warning("OpenAI API key not found")

if api_key_gemini:
    logger.info("Gemini API key loaded successfully")
else:
    logger.warning("Gemini API key not found")

client = OpenAI(api_key=api_key_openAI)
genai.configure(api_key=api_key_gemini)

class AIRequestType(Enum):
    EXECUTIVE_SUMMARY = "executive_summary"
    MARKET_ENABLERS = "market_enablers"
    INDUSTRY_EXPANSION = "industry_expansion"
    INDUSTRY_EXPANSION_1 = "industry_expansion_1" 
    INVESTMENT_CHALLENGES = "investment_challenges"
    COMPANY_INFO = "company_info"
    RESEARCH_JOURNALS = "research_journals"
    INDUSTRY_ASSOCIATIONS = "industry_associations"

class AIService:
    def __init__(self, openai_client, gemini_api_key):
        logger.info("Initializing AI Service")
        self.openai_client = openai_client
        self.gemini_configured = False
        if gemini_api_key:
            genai.configure(api_key=gemini_api_key)
            self.gemini_configured = True
            logger.info("Gemini configured successfully")
        else:
            logger.warning("Gemini not configured")
    
    def generate_content(self, request_type: AIRequestType, context: Dict[str, Any], existing_title: str = None) -> Any:
        start_time = time.time()
        logger.info(f"Generating content for: {request_type.value}")
        
        try:
            if request_type == AIRequestType.EXECUTIVE_SUMMARY:
                result = self._generate_executive_summary(context)
            elif request_type == AIRequestType.MARKET_ENABLERS:
                result = self._generate_market_enablers(context)
            elif request_type == AIRequestType.INDUSTRY_EXPANSION:
                result = self._generate_industry_expansion(context)
            elif request_type == AIRequestType.INDUSTRY_EXPANSION_1: 
                result = self._generate_industry_expansion_1(context, existing_title)
            elif request_type == AIRequestType.INVESTMENT_CHALLENGES:
                result = self._generate_investment_challenges(context)
            elif request_type == AIRequestType.COMPANY_INFO:
                result = self._generate_company_info(context)
            elif request_type == AIRequestType.RESEARCH_JOURNALS:
                result = self._generate_research_journals(context)
            elif request_type == AIRequestType.INDUSTRY_ASSOCIATIONS:
                result = self._generate_industry_associations(context)
            
            elapsed = time.time() - start_time
            timing_logger.info(f"{request_type.value} completed in {elapsed:.2f}s")
            logger.info(f"Content generation completed for: {request_type.value}")
            return result
            
        except Exception as e:
            elapsed = time.time() - start_time
            logger.error(f"Error generating {request_type.value} after {elapsed:.2f}s: {str(e)}")
            raise

    def generate_content_parallel(self, ai_context: Dict[str, Any]) -> Dict[str, Any]:
        start_time = time.time()
        logger.info("Starting parallel AI content generation")
        results = {}
        
        phase1_tasks = {
            'executive_summary': (AIRequestType.EXECUTIVE_SUMMARY, ai_context),
            'market_enablers': (AIRequestType.MARKET_ENABLERS, ai_context),
            'industry_expansion': (AIRequestType.INDUSTRY_EXPANSION, ai_context),
            'investment_challenges': (AIRequestType.INVESTMENT_CHALLENGES, ai_context),
            'research_journals': (AIRequestType.RESEARCH_JOURNALS, ai_context),
            'industry_associations': (AIRequestType.INDUSTRY_ASSOCIATIONS, ai_context),
            'company_info': (AIRequestType.COMPANY_INFO, ai_context)
        }
        
        logger.info(f"Phase 1: Executing {len(phase1_tasks)} tasks in parallel")
        
        with concurrent.futures.ThreadPoolExecutor(max_workers=7) as executor:
            future_to_key = {
                executor.submit(self.generate_content, request_type, context): key
                for key, (request_type, context) in phase1_tasks.items()
            }
            
            for future in concurrent.futures.as_completed(future_to_key):
                key = future_to_key[future]
                try:
                    results[key] = future.result()
                    logger.info(f"Phase 1 task completed: {key}")
                except Exception as exc:
                    logger.error(f'Phase 1 task {key} generated an exception: {exc}')
                    raise exc
        
        logger.info("Phase 1 completed, starting Phase 2")
        industry_title = results['industry_expansion']['title']
        results['industry_expansion_1'] = self.generate_content(
            AIRequestType.INDUSTRY_EXPANSION_1, 
            ai_context, 
            industry_title
        )
        
        elapsed = time.time() - start_time
        timing_logger.info(f"Parallel AI generation completed in {elapsed:.2f}s")
        logger.info("Parallel AI content generation completed successfully")
        return results
    
    def _generate_executive_summary(self, context: Dict[str, Any]) -> str:
        logger.info("Generating executive summary")
        first_line = (f"The {context['headline']} is valued at {context['cur']} {context['rev_current']} "
                    f"{context['value_in']} in {context['base_year']}, and is expected to reach "
                    f"{context['cur']} {context['rev_future']} {context['value_in']} by {context['forecast_year']}. "
                    f"The market shows a steady CAGR of {context.get('cagr')}% from 2025 to 2032.")
        
        prompt = f"Write an executive summary for {context['headline']} focusing on key market drivers, trends, and growth factors within 50 words stricly. Do not include market size or revenue figures as they are already provided. Focus on qualitative insights about market dynamics, key players, and future outlook. ( start directly from setence without any intro like 'The executive summary is...')"
        
        response = self.openai_client.chat.completions.create(
            model="gpt-5-mini",
            messages=[{"role": "user", "content": prompt}]
        )
        ai_summary = response.choices[0].message.content
        full_summary = f"{first_line} {ai_summary}"
        logger.info("Executive summary generated successfully")
        return full_summary
    
    def _generate_market_enablers(self, context: Dict[str, Any]) -> str:
        logger.info("Generating market enablers")
        prompt = f'Write an executive summary about key market enablers (2 points) for {context["headline"]}, each 50 words strickly. Return a Python list like ["heading: context", "heading: context"].'
        response = self.openai_client.chat.completions.create(
            model="gpt-5-mini",
            messages=[{"role": "user", "content": prompt}]
        )
        result = "\n".join(ast.literal_eval(response.choices[0].message.content))
        logger.info("Market enablers generated successfully")
        return result
    
    def _generate_industry_expansion(self, context: Dict[str, Any]) -> Dict[str, Any]:
        logger.info("Generating industry expansion")
        prompt = (
            f'Write one TOP Key Driver for the {context["headline"]} market. '
            f'Include a clear heading for the driver. '
            f'Return the output strictly as a Python dictionary with the following structure: '
            f'{{"title": "7–10 words", "paragraphs": ["paragraph1", "paragraph2", "paragraph3","parapragh4"]}}. '
            f'Each paragraph should be 80  words strict, qualitative in tone give 4 paragraphs, '
            f'and include real-world examples and facts. '
            f'Do not include market size, numbers, or links.'
        ) 

        response = self.openai_client.chat.completions.create(
            model="gpt-5-mini",
            messages=[{"role": "user", "content": prompt}]
        )

        result = ast.literal_eval(response.choices[0].message.content)
        logger.info("Industry expansion generated successfully")
        return result
    
    def _generate_industry_expansion_1(self, context: Dict[str, Any], existing_title: str = None) -> Dict[str, Any]:
        logger.info(f"Generating industry expansion 1 (avoiding title: {existing_title})")
        existing_title_instruction = ""
        if existing_title:
            existing_title_instruction = f' Do not use "{existing_title}" as the title - generate a completely different driver.'
        
        prompt = (
            f'Write one TOP Key Driver for the {context["headline"]} market that is DIFFERENT from previous drivers.{existing_title_instruction} '
            f'Include a clear heading for the driver. '
            f'Return the output strictly as a Python dictionary with the following structure: '
            f'{{"title": "7–10 words", "paragraphs": ["paragraph1", "paragraph2", "paragraph3","parapragh4"]}}. '
            f'Each paragraph should be 80 words strict, qualitative in tone ,give 4 paragraphs must, '
            f'and include real-world examples and facts. '
            f'Do not include market size, numbers, or links. '
            f'Focus on a unique aspect not covered by other drivers.'
        )

        response = self.openai_client.chat.completions.create(
            model="gpt-5-mini",
            messages=[{"role": "user", "content": prompt}]
        )

        result = ast.literal_eval(response.choices[0].message.content)
        logger.info("Industry expansion 1 generated successfully")
        return result
    
    def _generate_investment_challenges(self, context: Dict[str, Any]) -> str:
        logger.info("Generating investment challenges")
        prompt = (
            f'Write one TOP Key MARKET RESTRAINTS or CHALLENGES for the {context["headline"]} market. '
            f'Include a clear heading for the driver. '
            f'Return the output strictly as a Python dictionary with the following structure: '
            f'{{"title": "7–10 words", "paragraphs": ["paragraph1", "paragraph2", "paragraph3","parapragh4"]}}. '
            f'Each paragraph should be 80 words strict, qualitative in tone, give 4 paragraphs must'
            f'and include real-world examples and facts. '
            f'Do not include market size, numbers, or links.'
        )

        response = self.openai_client.chat.completions.create(
            model="gpt-5-mini",
            messages=[{"role": "user", "content": prompt}]
        )

        result = ast.literal_eval(response.choices[0].message.content)
        logger.info("Investment challenges generated successfully")
        return result
    
    def _generate_company_info(self, context: Dict[str, Any]) -> Dict[str, str]:
        logger.info(f"Generating company info for: {context['company_name']}")
        prompt = f'''Generate information about {context["company_name"]} in the "{context["headline"]}" domain. 
        Return the information in the following JSON format:
        {{
            "company_name": "{context["company_name"]}",
            "headquarters": "",
            "employee_count": "",
            "revenue": "",
            "top_product": "",
            "description_product": "",
            "estd": "",
            "website": "",
            "geographic_presence": "",
            "ownership": "",
            "short_description_company": ""
        }}
        geographic_presence only choose between from  Global, North America, Europe, Asia Pacific, Latin America, Middle East & Africa
        The short_description_company should be around 100 words. I want you to act as a Research Analyst and give Company Overview of "{context["company_name"]}" in around 10-11 lines (In one paragraph only) which should not talk about Headquarter Country, Establishment/Foundation Year, Number of Employees or Revenue and should not use any marketing/promotional words like, largest, prominent, diversified, recognized, among others (You can talk about its product/service related to {context["headline"]}, market presence, business strategy, recent developments, etc) like this for tone:
        Schlumberger Ltd (SLB) provides technology for reservoir characterization, production, drilling and processing to the oil and gas industry. The company supplies its products and services to the industry, from exploration through production and integrated pipeline solutions for hydrocarbon recovery. SLB's products and services include open-hole and cased-hole wireline logging; drilling services; well completion services, including well testing and artificial lift; well services such as cementing, coiled tubing, stimulations, and sand control; interpretation and consulting services; and integrated project management. The company has an operational presence in North America, Latin America, Europe and Africa, the Middle East and Asia. SLB is headquartered in Houston, Texas, the US..
.       website should be the official website no Https ot http.
        revenue should be in the format " X.XX billion" or " X.XX million" and should be correct 2024 data in USD only correct data must strict.
        ownership should be either "Public" or "Private".
        top product should be a product or service relevant to the headline market.
        description_product should be 50 words describing the top product.
        estd is year of establishment should be correct data.
        headquarters should be "Country" format and should be correct data.
        employee_count should be in "X,XXX" or "XX,XXX" format and should be correct data.
        Return ONLY valid JSON, no additional text. no urls/citations for references.'''
        
        response = client.responses.create(
            model="gpt-5",
            tools=[{
                "type": "web_search_preview",
                "search_context_size": "medium",
            }],
            input=[
                {"role": "system", "content": "You are a JSON generator. Always return valid JSON and nothing else."},
                {"role": "user", "content": prompt}
            ]
        )          
        content = response.output_text.strip()
        
        if content.startswith("```json"):
            content = content[7:]
        if content.endswith("```"):
            content = content[:-3]
        content = content.strip()
        
        result = json.loads(content)
        logger.info("Company info generated successfully")
        return result

    def _generate_research_journals(self, context: Dict[str, Any]) -> List[str]:
        logger.info("Generating research journals")
        market_name = context.get('headline', 'Technology Market')
        
        response = self.openai_client.chat.completions.create(
            model="gpt-5-mini",
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are a JSON generator. Provide the names of research journals related to the specified market "
                        "in JSON format. Only include the names as strings, no additional information "
                        "is needed. Search established, reputable journals.\n\n"
                        "Give 5 journal names.\n\n"
                        "**Output format must be a JSON object with a 'journals' key containing an array of strings:**\n"
                        '{"journals": ["Journal Name 1", "Journal Name 2"]}\n'
                        "If there are no journals for the given market, return: {\"journals\": []}"
                    )
                },
                {
                    "role": "user",
                    "content": f"Find research journals for: {market_name}"
                }
            ],
            response_format={"type": "json_object"}
        )
        
        json_response = json.loads(response.choices[0].message.content)
        journals = json_response.get('journals', [])
        
        default_journals = [
            "Journal of Market Research",
            "International Business Review",
            "Strategic Management Journal",
            "Harvard Business Review",
            "Industrial Marketing Management"
        ]
        
        if len(journals) < 5:
            journals.extend(default_journals[len(journals):5])
        
        logger.info(f"Research journals generated: {len(journals)} items")
        return journals[:5]  
    
    def _generate_industry_associations(self, context: Dict[str, Any]) -> List[str]:
        logger.info("Generating industry associations")
        market_name = context.get('headline', 'Technology Market')
        
        response = self.openai_client.chat.completions.create(
            model="gpt-5-mini",
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are a JSON generator. Provide the names of industry associations or government organizations "
                        "related to the specified market in JSON format. Only include the names "
                        "as strings, no additional information is needed. Search for highly relevant organizations "
                        "to the market name (exclude private company names). Give 5.\n\n"
                        "**Output format must be a JSON object with an 'associations' key containing an array of strings:**\n"
                        '{"associations": ["Association Name 1", "Association Name 2"]}\n'
                        "If there are no relevant associations or organizations for the given market, "
                        "return: {\"associations\": []}"
                    )
                },
                {
                    "role": "user",
                    "content": f"Find industry associations and government organizations for: {market_name}"
                }
            ],
            response_format={"type": "json_object"}
        )
        
        json_response = json.loads(response.choices[0].message.content)
        associations = json_response.get('associations', [])
        
        default_associations = [
            "Global Industry Alliance",
            "International Trade Association",
            "National Business Federation",
            "Industry Development Council",
            "Professional Standards Organization"
        ]
        
        if len(associations) < 5:
            associations.extend(default_associations[len(associations):5])
        
        logger.info(f"Industry associations generated: {len(associations)} items")
        return associations[:5]  

logger.info("Initializing AI Service globally")
ai_service = AIService(client, api_key_gemini)

class TaxonomyBoxGenerator:
    COLORS = {
        "purple": RGBColor(0x31, 0x09, 0x7E),
        "orange": RGBColor(255, 102, 51),
        "teal": RGBColor(0, 179, 152),
        "blue": RGBColor(0, 162, 232),
        "dark_blue": RGBColor(36, 64, 142),
        "white": RGBColor(255, 255, 255),
        "light_gray": RGBColor(0xF2, 0xF2, 0xF2),
        "text_dark": RGBColor(0, 0, 0),
        "new_blue": RGBColor(0x00, 0x70, 0xC0),     
        "light_green": RGBColor(0x92, 0xD0, 0x50),   
        "yellow_orange": RGBColor(0xFF, 0xC0, 0x00),
        "dark_red": RGBColor(0xC0, 0x00, 0x00),     
        "rose": RGBColor(0xF8, 0x78, 0x84),         
        "light_black": RGBColor(0x7F, 0x7F, 0x7F), 
        "dark_teal": RGBColor(0x00, 0xA8, 0x8F),   
        "turquoise": RGBColor(0x33, 0xC5, 0xF0),    
        "new_purple": RGBColor(0x59, 0x46, 0x8F),   
    }
    
    BOX_HEADER_COLORS = [
        COLORS["new_blue"],      
        COLORS["light_green"],   
        COLORS["yellow_orange"], 
        COLORS["dark_red"],      
        COLORS["rose"],          
        COLORS["light_black"],   
        COLORS["dark_teal"],     
        COLORS["turquoise"],     
        COLORS["new_purple"],    
    ]

    def __init__(self, presentation):
        self.prs = presentation
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height
        self.left_margin, self.top_margin, self.right_margin, self.bottom_margin = (
            Inches(0.5),
            Inches(2),
            Inches(0.5),
            Inches(0.8),
        )
        self.h_spacing, self.v_spacing = Inches(0.2), Inches(0.2)

    def _add_category_box(
        self, slide, category, content, left, top, max_width, max_height, color_index
    ):
        header_color = self.BOX_HEADER_COLORS[color_index % len(self.BOX_HEADER_COLORS)]
        
        header_height = Inches(0.3)
        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, max_width, header_height
        )
        header.fill.solid()
        header.fill.fore_color.rgb = header_color
        header.line.color.rgb = header_color
        p = header.text_frame.paragraphs[0]
        p.text = category
        p.font.size, p.font.bold, p.font.color.rgb, p.alignment = (
            Pt(11),
            True,
            self.COLORS["white"],
            PP_ALIGN.CENTER,
        )

        content_box_height = max_height - header_height + Inches(0.2)
        content_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left,
            top + header_height,
            max_width,
            content_box_height,
        )
        content_box.fill.solid()
        content_box.fill.fore_color.rgb = self.COLORS["light_gray"]
        content_box.line.color.rgb = self.COLORS["light_gray"]
        tf = content_box.text_frame
        tf.word_wrap, tf.vertical_anchor = True, MSO_VERTICAL_ANCHOR.TOP
        self._add_list_content(tf, content)

    def _add_list_content(self, text_frame, content):
        text_frame.margin_bottom = Pt(12)
        if text_frame.paragraphs:
            text_frame.paragraphs[0].text = ""

        p = text_frame.paragraphs[0]
        pPr = p._p.get_or_add_pPr()

        lst = pPr.find(qn("a:lstStyle"))
        if lst is None:
            lst = OxmlElement("a:lstStyle")
            pPr.append(lst)

        for i, item in enumerate(content):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            p.text = item
            p.alignment = PP_ALIGN.LEFT

            pPr = p._p.get_or_add_pPr()
            pPr.set("marL", str(int(Pt(15).emu)))
            pPr.set("indent", str(int(Pt(-15).emu)))

            marL = OxmlElement("a:marL")
            marL.set("val", str(int(Pt(50).emu)))
            pPr.append(marL)

            indent = OxmlElement("a:indent")
            indent.set("val", str(int(Pt(-19).emu)))
            pPr.append(indent)

            buChar = OxmlElement("a:buChar")
            buChar.set("char", "○")
            pPr.append(buChar)

            buFont = OxmlElement("a:buFont")
            buFont.set("typeface", "Symbol")
            pPr.append(buFont)

            buClr = OxmlElement("a:buClr")
            srgbClr = OxmlElement("a:srgbClr")
            srgbClr.set("val", "000000")
            buClr.append(srgbClr)
            pPr.append(buClr)

            for run in p.runs:
                run.font.size = Pt(10)
                run.font.color.rgb = self.COLORS["text_dark"]

            if not pPr.find(qn("a:spcAft")):
                spcAft = OxmlElement("a:spcAft")
                spcVal = OxmlElement("a:spcPts")
                spcVal.set("val", "600")
                spcAft.append(spcVal)
                pPr.append(spcAft)

            if not pPr.find(qn("a:lnSpc")):
                lnSpc = OxmlElement("a:lnSpc")
                spcPts = OxmlElement("a:spcPts")
                spcPts.set("val", "1800")
                lnSpc.append(spcPts)
                pPr.append(lnSpc)

    def add_taxonomy_boxes(self, slide_index, taxonomy_data):
        logger.info(f"Adding taxonomy boxes to slide {slide_index}")
        slide = self.prs.slides[slide_index]
        available_width = self.slide_width - self.left_margin - self.right_margin
        num_categories = len(taxonomy_data)
        boxes_per_row = min(5, num_categories)
        box_width = (
            available_width - (boxes_per_row - 1) * self.h_spacing
        ) / boxes_per_row

        rows, current_row, current_row_width = [], [], 0
        color_index = 0         
        for category, hierarchy in taxonomy_data.items():
            item_count = len(hierarchy)
            box_height = max(
                Inches(1), Inches(0.43) + (item_count * Inches(0.17) * 1.2)
            )
            if current_row_width + box_width > available_width and current_row:
                rows.append(current_row)
                current_row, current_row_width = [], 0
            current_row.append(
                {
                    "category": category, 
                    "content": hierarchy, 
                    "height": box_height,
                    "color_index": color_index  
                }
            )
            current_row_width += box_width + self.h_spacing
            color_index += 1  
            
        if current_row:
            rows.append(current_row)

        current_top = self.top_margin
        for row in rows:
            row_max_height = max(box["height"] for box in row)
            row_width = len(row) * box_width + (len(row) - 1) * self.h_spacing
            left_start = self.left_margin + (available_width - row_width) / 2
            for i, box in enumerate(row):
                left = left_start + i * (box_width + self.h_spacing)
                self._add_category_box(
                    slide,
                    box["category"],
                    box["content"],
                    left,
                    current_top,
                    box_width,
                    box["height"],
                    box["color_index"]  
                )
            current_top += row_max_height + self.v_spacing
        
        logger.info(f"Taxonomy boxes added successfully to slide {slide_index}")


def replace_text_in_presentation(prs, slide_data_dict):
    logger.info("Starting text replacement in presentation")
    for slide_idx, slide in enumerate(prs.slides):
        data = slide_data_dict.get(slide_idx, {})
        if not data:
            continue
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                for key, value in data.items():
                    token = f"{{{{{key}}}}}"
                    if token in p.text:
                        inline_text = p.text
                        p.text = inline_text.replace(token, str(value))
    logger.info("Text replacement in presentation completed")


def replace_text_in_tables(prs, slide_indices, slide_data_dict):
    logger.info(f"Starting text replacement in tables for slides: {slide_indices}")
    for idx in slide_indices:
        if idx >= len(prs.slides):
            continue
        slide = prs.slides[idx]
        data = slide_data_dict.get(idx, {})
        if not data:
            continue
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            for row in shape.table.rows:
                for cell in row.cells:
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            for key, value in data.items():
                                if f"{{{{{key}}}}}" in run.text:
                                    run.text = run.text.replace(
                                        f"{{{{{key}}}}}", str(value)
                                    )
    logger.info("Text replacement in tables completed")


def get_rgb_color_safe(font):
    try:
        return font.color.rgb
    except AttributeError:
        return None


def replace_text_preserving_color(paragraph, placeholder, new_text):
    full_text = "".join(run.text for run in paragraph.runs)

    if placeholder not in full_text:
        return

    for run in paragraph.runs:
        if placeholder in run.text:
            font_color = get_rgb_color_safe(run.font)
            run.text = run.text.replace(placeholder, new_text)
            if font_color:
                run.font.color.rgb = font_color
            break


def replace_text_in_paragraph(paragraph, placeholder, new_text):
    full_text = "".join(run.text for run in paragraph.runs)
    if placeholder in full_text:
        replacement_text = str(new_text) if new_text is not None else ""
        
        if not replacement_text.strip():
            new_full_text = full_text.replace(placeholder, "").strip()
            import re
            new_full_text = re.sub(r'\n\s*\n', '\n', new_full_text)
        else:
            new_full_text = full_text.replace(placeholder, replacement_text)
        
        for run in paragraph.runs:
            run.text = ""
        if paragraph.runs:
            paragraph.runs[0].text = new_full_text
        else:
            paragraph.add_run().text = new_full_text


def set_cell_border(cell):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    for line_dir in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
        ln = OxmlElement(line_dir)
        ln.set("w", "12700")

        solidFill = OxmlElement("a:solidFill")
        srgbClr = OxmlElement("a:srgbClr")
        srgbClr.set("val", "A6A6A6")
        solidFill.append(srgbClr)
        ln.append(solidFill)

        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        ln.set("algn", "ctr")

        tcPr.append(ln)


def validate_segment_hierarchy(segment_text):
    logger.info("Validating segment hierarchy")
    lines = segment_text.strip().split('\n')
    errors = []
    last_main_number = 0
    last_sub_numbers = {}

    for i, line in enumerate(lines):
        line = line.strip()
        if not line:
            continue

        match = re.match(r'^(\d+(?:\.\d+)*)\.\s*(.+)$', line)
        if not match:
            errors.append(f"Line {i + 1}: Invalid format")
            continue

        number_parts = [int(n) for n in match.group(1).split('.')]
        depth = len(number_parts)

        if depth == 1:
            if number_parts[0] != last_main_number + 1:
                errors.append(f"Line {i + 1}: Expected main number {last_main_number + 1}, got {number_parts[0]}")
            last_main_number = number_parts[0]  
            last_sub_numbers = {}
        elif depth == 2:
            main_num = number_parts[0]
            sub_num = number_parts[1]
            
            if main_num != last_main_number:
                errors.append(f"Line {i + 1}: Sub-item doesn't match current main number")
            
            expected_sub = last_sub_numbers.get(main_num, 0) + 1
            if sub_num != expected_sub:
                errors.append(f"Line {i + 1}: Expected sub-number {main_num}.{expected_sub}")
            last_sub_numbers[main_num] = sub_num
        elif depth == 3:
            main_num = number_parts[0]
            sub_num = number_parts[1]
            sub_sub_num = number_parts[2] 
            
            if main_num != last_main_number:
                errors.append(f"Line {i + 1}: Sub-sub-item doesn't match current main number")
            
            key = f"{main_num}.{sub_num}"
            expected_sub_sub = last_sub_numbers.get(key, 0) + 1
            if sub_sub_num != expected_sub_sub:
                errors.append(f"Line {i + 1}: Expected sub-sub-number {main_num}.{sub_num}.{expected_sub_sub}")
            last_sub_numbers[key] = sub_sub_num

    if errors:
        logger.warning(f"Segment hierarchy validation found {len(errors)} errors")
    else:
        logger.info("Segment hierarchy validation passed")
    return errors


def generate_actual_data():
    data = [
        [2019, 7.0, 5.0, 4.0, 3.5, 3.1, 2.5, 2.0],
        [2020, 7.5, 5.4, 4.3, 3.7, 3.3, 2.7, 2.2],
        [2021, 8.1, 5.7, 4.6, 4.0, 3.6, 2.9, 2.3],
        [2022, 8.7, 6.1, 5.0, 4.3, 3.8, 3.1, 2.5],
        [2023, 9.3, 6.6, 5.3, 4.7, 4.1, 3.3, 2.7],
        [2024, 9.9, 7.1, 5.8, 5.0, 4.4, 3.6, 2.9],
        [2025, 10.7, 7.7, 6.2, 5.4, 4.8, 3.8, 3.1],
        [2026, 11.6, 8.2, 6.7, 5.8, 5.1, 4.1, 3.3],
        [2027, 12.4, 8.9, 7.2, 6.2, 5.5, 4.5, 3.6],
        [2028, 13.4, 9.6, 7.7, 6.7, 5.9, 4.8, 3.8],
        [2029, 14.4, 10.3, 8.4, 7.2, 6.4, 5.1, 4.2],
        [2030, 15.4, 11.1, 9.0, 7.7, 6.8, 5.6, 4.5],
        [2031, 16.6, 11.9, 9.7, 8.4, 7.3, 5.9, 4.8],
        [2032, 17.8, 12.8, 10.4, 9.0, 7.9, 6.4, 5.2]
    ]
    return data


def parse_segment_input(segment_input: str) -> Dict[str, Dict]:
    logger.info("Parsing segment input")
    lines = segment_input.strip().split("\n")
    nested_dict = {}
    level_stack = []
    for line in lines:
        if not line.strip():
            continue
        if ". " in line:
            key, value = line.split(". ", 1)
        elif "." in line:
            parts = line.split(".")
            for i in range(len(parts) - 1, 0, -1):
                try:
                    number_part = ".".join(parts[:i])
                    text_part = ".".join(parts[i:])
                    [int(n) for n in number_part.split('.')]
                    key = number_part
                    value = text_part
                    break
                except ValueError:
                    continue
        else:
            continue
            
        parts = key.split(".")
        depth = len(parts)
        label = value.strip()
        level_stack = level_stack[: depth - 1]
        current = nested_dict
        for k in level_stack:
            current = current[k]
        current[label] = {}
        level_stack.append(label)
    logger.info(f"Segment input parsed successfully: {len(nested_dict)} main categories")
    return nested_dict

def generate_toc_data(nested_dict: Dict, headline: str, forecast_period: str, user_segment: str, kmi_items: List[str] = None) -> Dict[str, int]:
    logger.info("Generating Table of Contents data")
    toc_start_levels = {
        "1. Introduction": 0,
        "1.1. Objectives of the Study": 1,
        "1.2. Market Definition & Scope": 1,
        "2. Research Methodology": 0,
        "2.1. Research Process": 1,
        "2.2. Secondary & Primary Data Methods": 1,
        "2.3. Market Size Estimation Methods": 1,
        "2.4. Market Assumptions & Limitations": 1,
        "3. Executive Summary": 0,
        "3.1. Global Market Outlook": 1,
        "3.2. Key Market Highlights": 1,
        "3.3. Segmental Overview": 1,
        "4. Market Dynamics & Outlook": 0,
        "4.1. Macro-Economic Indicators​": 1,
        "4.2. Drivers & Opportunities": 1,
        "4.3. Restraints & Challenges": 1,
        "4.4. Supply Side Trends": 1,
        "4.5. Demand Side Trends": 1,
        "4.6. Porter's Analysis & Impact": 1,
        "4.6.1. Competitive Rivalry": 2,
        "4.6.2. Threat of substitutes": 2,
        "4.6.3. Bargaining power of buyers": 2,
        "4.6.4. Threat of new entrants": 2,
        "4.6.5. Bargaining power of suppliers": 2,
    }

    kmi_section = {"5. Key Market Insights": 0}

    default_kmi_items = [
        "Key Success Factors",
        "Market Impacting Factors", 
        "Top Investment Pockets",
        "Market Attractiveness Index, 2024",
        "Market Ecosystem",
        "PESTEL Analysis",
        "Pricing Analysis",
        "Regulatory Landscape",
      
    ]

    all_kmi_items = default_kmi_items.copy()
    if kmi_items:
        all_kmi_items.extend(kmi_items)

    for i, kmi_item in enumerate(all_kmi_items, start=1):
        kmi_section[f"5.{i}. {kmi_item}"] = 1
    toc_mid = {}
    main_index = 6
    for type_index, (type_name, points) in enumerate(nested_dict.items(), start=main_index):
        toc_mid[
            f"{type_index}. {headline} Size by {type_name} (2019-2032)"
        ] = 0
        point_count = 1
        for point, subpoints in points.items():
            toc_mid[f"{type_index}.{point_count}. {point}"] = 1
            if subpoints:
                for sp_count, sub in enumerate(subpoints, start=1):
                    toc_mid[f"{type_index}.{point_count}.{sp_count}. {sub}"] = 2
            point_count += 1
        


    x = len(list(nested_dict.keys())) + 6
    toc_end_levels = {
        f"{x}. Global {headline} Size by Region (2019-2032)": 0,
        f"{x}.1. North America ({user_segment})": 1,
        f"{x}.1.1. US": 2,
        f"{x}.1.2. Canada": 2,
        f"{x}.2. Europe ({user_segment})": 1,
        f"{x}.2.1. UK": 2,
        f"{x}.2.2. Germany": 2,
        f"{x}.2.3. Spain": 2,
        f"{x}.2.4. France": 2,
        f"{x}.2.5. Italy": 2,
        f"{x}.2.6. Rest of Europe": 2,
        f"{x}.3. Asia-Pacific ({user_segment})": 1,
        f"{x}.3.1. China": 2,
        f"{x}.3.2. India": 2,
        f"{x}.3.3. Japan": 2,
        f"{x}.3.4. South Korea": 2,
        f"{x}.3.5. Rest of Asia Pacific": 2,
        f"{x}.4. Latin America ({user_segment})": 1,
        f"{x}.4.1. Brazil": 2,
        f"{x}.4.2. Mexico": 2,
        f"{x}.4.3. Rest of Latin America": 2,
        f"{x}.5. Middle East & Africa ({user_segment})": 1,
        f"{x}.5.1. GCC Countries": 2,
        f"{x}.5.2. South Africa": 2,
        f"{x}.5.3. Rest of Middle East & Africa": 2,
        f"{x+1}. Competitive Landscape": 0,
        f"{x+1}.1. Competitive Dashboard": 1,
        f"{x+1}.2. Market Positioning of Key Players, 2024": 1,
        f"{x+1}.3. Strategies Adopted by Key Market Players": 1,
        f"{x+1}.4. Recent Developments in the Market": 1,
        f"{x+1}.5. Company Market Share Analysis, 2024": 1,
        f"{x+2}. Key Company Profiles": 0,
    }
    
    logger.info(f"TOC data generated with {len(toc_start_levels) + len(kmi_section) + len(toc_mid) + len(toc_end_levels)} items")
    return {**toc_start_levels, **kmi_section, **toc_mid, **toc_end_levels}


def add_toc_to_slides(prs: Presentation, toc_data_levels: Dict[str, int], toc_slide_indices: List[int]):
    logger.info(f"Adding TOC to slides: {toc_slide_indices}")
    for i in toc_slide_indices:
        slide = prs.slides[i]
        table_shape = slide.shapes.add_table(
            17, 2, Inches(2.8), Inches(0.5), Inches(10), Inches(6)
        )
        table = table_shape.table
        for row in table.rows:
            for cell in row.cells:
                cell.text = ""
                cell.fill.background()
                tcPr = cell._tc.get_or_add_tcPr()
                for border_tag in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
                    tcPr.append(
                        parse_xml(
                            f'<{border_tag} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:noFill/></{border_tag}>'
                        )
                    )

    content_items = list(toc_data_levels.keys())
    content_index = 0
    for i in toc_slide_indices:
        table = prs.slides[i].shapes[-1].table
        for col in range(2):
            for row in range(17):
                if content_index >= len(content_items):
                    break
                cell, key = table.cell(row, col), content_items[content_index]
                level = toc_data_levels[key]
                para = cell.text_frame.paragraphs[0]
                para.text = "          " * level + key
                font = para.font
                font.color.rgb, font.size, font.name = RGBColor(0, 0, 0), Pt(11), "Poppins"
                if key.startswith("The following companies") or key.startswith("Note :"):
                    font.size = Pt(9)
                    font.color.rgb, font.bold = RGBColor(112, 48, 160), True
                else:
                    font.size = Pt(11)
                if level == 0 and not key.startswith("The following companies") and not key.startswith("Note :"):
                    font.color.rgb, font.bold = RGBColor(112, 48, 160), True
                else:
                    font.color.rgb = RGBColor(0, 0, 0)
                    font.bold = False
                content_index += 1
    logger.info("TOC added to slides successfully")


def create_chart_on_slide(slide: Any, data: List[List], chart_columns: List[str], 
                         left: float, top: float, width: float, height: float):
    logger.info(f"Creating chart with {len(chart_columns)} series and {len(data)} data points")
    chart_data = CategoryChartData()
    chart_data.categories = [str(row[0]) for row in data]

    num_series = min(len(chart_columns), 7)
    for i in range(num_series):
        chart_data.add_series(chart_columns[i], [row[i + 1] for row in data])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, left, top, width, height, chart_data
    ).chart
    
    chart.plots[0].gap_width = 150  
    chart.chart_style = 2
    chart.has_title = False

    chart.has_legend = True
    chart.legend.font.size = Pt(8)  
    chart.legend.font.name = "Poppins"
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    value_axis = chart.value_axis
    value_axis.visible = False
    value_axis.has_major_gridlines = False
    value_axis.has_minor_gridlines = False
    value_axis.tick_label_position = XL_TICK_LABEL_POSITION.NONE
    value_axis.major_tick_mark = XL_TICK_MARK.NONE
    value_axis.minor_tick_mark = XL_TICK_MARK.NONE

    category_axis = chart.category_axis
    category_axis.major_tick_mark = XL_TICK_MARK.NONE
    category_axis.minor_tick_mark = XL_TICK_MARK.NONE
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(10)
    cat_axis.tick_labels.font.name = "Poppins"
    cat_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
    logger.info("Chart created successfully")


def clean_filename(filename):
    logger.info(f"Cleaning filename: {filename}")
    invalid_chars = '<>:"/\\|?*'
    for char in invalid_chars:
        filename = filename.replace(char, '_')
    
    filename = ' '.join(filename.split())  
    filename = filename[:100]  
    
    logger.info(f"Cleaned filename: {filename}")
    return filename


@app.route('/')
def index():
    logger.info("Index page accessed")
    return render_template('index.html')


@app.route('/generate-ppt', methods=['POST'])
def generate_ppt():
    start_time = time.time()
    request_id = f"req_{int(time.time())}"
    logger.info(f"[{request_id}] PPT generation request started")
    
    try:
        form_data = request.form
        required_fields = [
            'headline', 'headline_2', 'historical_year', 'base_year',
            'forecast_year', 'forecast_period', 'cur', 'value_in',
            'rev_current', 'rev_future', 'segment_input', 'companies','cagr',
        ]
        
        missing_fields = []
        for field in required_fields:
            if not form_data.get(field, '').strip():
                missing_fields.append(field)
        
        if missing_fields:
            logger.warning(f"[{request_id}] Missing required fields: {missing_fields}")
            return jsonify({
                'error': 'Missing required fields',
                'fields': missing_fields
            }), 400
        
        segment_errors = validate_segment_hierarchy(form_data['segment_input'])
        if segment_errors:
            logger.warning(f"[{request_id}] Segment hierarchy validation failed: {len(segment_errors)} errors")
            return jsonify({
                'error': 'Invalid segment hierarchy',
                'details': segment_errors
            }), 400
        
        logger.info(f"[{request_id}] Form validation passed")
        
        headline = form_data['headline']
        headline_2 = headline.upper()
        headline_3 = headline_2.replace("GLOBAL", "").strip()
        historical_year = "2019-2023"
        base_year = "2024"
        forecast_year = "2032"
        forecast_period = "2025-2032"
        cur = "USD"
        value_in = form_data['value_in']
        currency = f"{cur} {value_in}"
        rev_current = form_data['rev_current']
        rev_future = form_data['rev_future']
        cagr = form_data.get('cagr')
        segment_input = form_data['segment_input']
        kmi_items = []
        kmi_input = form_data.get('kmi_items', '').strip()
        
        if kmi_input:
            kmi_items = [item.strip() for item in kmi_input.split('\n') if item.strip()]
            logger.info(f"[{request_id}] Custom KMI items provided: {len(kmi_items)}")
      
        def format_as_bullets(items_list):
            if not items_list:
                return ""
            return '\n'.join([f"{item}" for item in items_list])
        default_kmiitems = [
                    "Key Success Factors",
                    "Market Impacting Factors", 
                    "Top Investment Pockets",
                    "Market Attractiveness Index, 2024",
                    "Market Ecosystem",
                    "PESTEL Analysis",
                    "Pricing Analysis",
                    "Regulatory Landscape",
                ]
        default_kmi_bullets = format_as_bullets(default_kmiitems)
        user_kmi_bullets = format_as_bullets(kmi_items) if kmi_items else ""
        
        companies_input = form_data['companies'].strip()
        company_list = [company.strip() for company in companies_input.split('\n') if company.strip()]
        
        if not company_list:
            logger.warning(f"[{request_id}] No companies provided")
            return jsonify({
                'error': 'At least one company must be provided',
                'message': 'Please provide company names, one per line'
            }), 400
        
        logger.info(f"[{request_id}] Processing {len(company_list)} companies")
        
        nested_dict = parse_segment_input(segment_input)
        main_topic = list(nested_dict.keys())
        s_segment = "By " + "\nBy ".join(main_topic)
        user_segment = "By " + ", By ".join(main_topic)

        output_lines = []
        for main_type, points in nested_dict.items():
            line_parts = []
            for point, subpoints in points.items():
                if subpoints:
                    subpoint_str = ", ".join(subpoints.keys())
                    line_parts.append(f"{point} ({subpoint_str})")
                else:
                    line_parts.append(point)
            output_lines.append(f"By {main_type}: {', '.join(line_parts)}")
        output_lines.append(
            "By Region: North America, Europe, Asia-Pacific, Latin America, Middle East & Africa"
        )
        context = "\n".join(output_lines)
        logger.info(f"[{request_id}] Context generated successfully")
        
        toc_data_levels = generate_toc_data(nested_dict, headline, forecast_period, user_segment, kmi_items)

        ai_context = {
            'headline': headline,
            'value_in': value_in,
            'cur': cur,
            'historical_year': historical_year,
            'forecast_year': forecast_year,
            'base_year': base_year,
            'rev_current': rev_current,
            'rev_future': rev_future,
            'main_topic': main_topic[0] if main_topic else "Type 1",
            'currency': currency.upper(),
            'cagr': cagr,
            'company_name': company_list[0] 
        }
        
        logger.info(f"[{request_id}] Starting AI content generation")
        ai_start_time = time.time()
        ai_results = ai_service.generate_content_parallel(ai_context)
        ai_elapsed = time.time() - ai_start_time
        timing_logger.info(f"[{request_id}] AI content generation completed in {ai_elapsed:.2f}s")
        
        mpara_11 = ai_results['executive_summary']
        para_11 = ai_results['market_enablers']
        para_14_dict = ai_results['industry_expansion']
        industry_title = para_14_dict['title']
        para_14_dict_1 = ai_results['industry_expansion_1']
        para_15_dict = ai_results['investment_challenges']
        research_journals = ai_results['research_journals']
        industry_associations = ai_results['industry_associations']
        company_info = ai_results['company_info']
        
        industry_title_1 = para_15_dict['title'] 
        industry_title_2 = para_14_dict_1['title']
        para_14_1= '\n'.join(para_14_dict_1['paragraphs'])
        para_15 = '\n'.join(para_15_dict['paragraphs'])
        para_14 = '\n'.join(para_14_dict['paragraphs'])
        
        logger.info(f"[{request_id}] AI content extracted successfully")
        
        x = len(main_topic) + 6

        first_company = company_list[0]
        toc_data_levels[f"{x+2}.1. {first_company}"] = 1
        toc_data_levels[f"{x+2}.1.1. Company Overview"] = 2
        toc_data_levels[f"{x+2}.1.2. Product Portfolio Overview"] = 2
        toc_data_levels[f"{x+2}.1.3. Financial Overview"] = 2
        toc_data_levels[f"{x+2}.1.4. Key Developments"] = 2

        toc_data_levels["The following companies are listed for indicative purposes only. Similar information will be provided for each, with detailed financial data available exclusively for publicly listed companies."] = 0
        for i, name in enumerate(company_list[1:], start=2):
            toc_data_levels[f"{x+2}.{i}. {name}"] = 1
        toc_data_levels["Note : The list of companies mentioned are for indication purpose and are subject to change over the due course of the research"]=0

        toc_data_levels[f"{x+3}. Conclusion & Recommendation"] = 0
        table_taxonomy = {
            f"BY {key.upper()}": list(value.keys()) for key, value in nested_dict.items()
        }
       
        slide_data = {
            0: {
                "heading": headline_2,
                "timeline": f"HISTORIC YEAR {historical_year} AND FORECAST TO {forecast_year}",
                "context": context,
            },
            1: {
                "heading_2": f"{headline_2} ({currency.upper()})",
                "hyear": f"Historical Year - {historical_year}",
                "fyear": f"Forecast Year - {forecast_period}",
                "byear": f"Base Year - {base_year}",
            }, 
            2: {
                "heading_2": f"{headline_2} ({currency.upper()})",
                "hyear": f"Historical Year - {historical_year}",
                "fyear": f"Forecast Year - {forecast_period}",
                "byear": f"Base Year - {base_year}",
            },
            3: {
                "heading": headline_2,
                "timeline": f"HISTORIC YEAR {historical_year} FORECAST TO {forecast_year}",
            },
            8: {
                "heading": headline_2,
                "timeline": f"HISTORIC YEAR {historical_year} FORECAST TO {forecast_year}",
            },
            10: {
                "org_1": industry_associations[0] if len(industry_associations) > 0 else "Global Industry Alliance",
                "org_2": industry_associations[1] if len(industry_associations) > 1 else "International Trade Association",
                "org_3": industry_associations[2] if len(industry_associations) > 2 else "National Business Federation",
                "org_4": industry_associations[3] if len(industry_associations) > 3 else "Industry Development Council",
                "org_5": industry_associations[4] if len(industry_associations) > 4 else "Professional Standards Organization",
                "paper_1": research_journals[0] if len(research_journals) > 0 else "Journal of Market Research",
                "paper_2": research_journals[1] if len(research_journals) > 1 else "International Business Review",
                "paper_3": research_journals[2] if len(research_journals) > 2 else "Strategic Management Journal",
                "paper_4": research_journals[3] if len(research_journals) > 3 else "Harvard Business Review",
                "paper_5": research_journals[4] if len(research_journals) > 4 else "Industrial Marketing Management",
            },
            12: {
                "heading": headline_2,
                "timeline": f"HISTORIC YEAR {historical_year} FORECAST TO {forecast_year}",
            },
            13: {
                "heading_2": f"{headline_2} SIZE, ({currency.upper()})",
                "mpara": mpara_11,
                "para": para_11,
                "amount_1": rev_current,    
                "amount_2": rev_future,
            },
            14: {
                "heading": headline_2,
                "amount_1": f"{cur} {rev_current} {value_in.upper()} ",
                "amount_2": f"{rev_future} {value_in.upper()} {cur}",
            },
            15:  {
                "heading": headline_2,
                "timeline": f"HISTORIC YEAR {historical_year} FORECAST TO {forecast_year}",
            },
            17: {"industry_title": industry_title, "para": para_14},
            18: {"industry_title": industry_title_2, "para": para_14_1},

            19: {"industry_title": industry_title_1, "para": para_15},
            21:  {
                "heading": headline_2,
                "timeline": f"HISTORIC YEAR {historical_year} FORECAST TO {forecast_year}",
                "default_kmi": default_kmi_bullets,
                "user_kmi": user_kmi_bullets,
            },
            23: {
                "heading": headline_2,  
                "timeline": f"HISTORIC YEAR {historical_year} FORECAST TO {forecast_year}",
                "types": s_segment,
            },
            24: {
                "heading":  headline_2,
                "type_1": main_topic[0].upper() if main_topic else "Type 1",
                "timeline": "2019-2032",
                "cur": f"{cur.upper()} {value_in.upper()}",
            },
            25: {
                "heading": headline_2,
                "timeline": f"HISTORIC YEAR {historical_year} FORECAST TO {forecast_year}",
            },
            26: {"heading": headline_2, "timeline": "2019-2032", "cur": f"{cur.upper()} {value_in.upper()}"},
            27: {
                "2_heading": headline_3.upper(),
                "type_1": main_topic[0].upper() if main_topic else "Type 1",
                "timeline": "2019-2032",
                "cur": f"{cur.upper()} {value_in.upper()}",
            },
            28: {
                "2_heading": headline_3.upper(),
                "type_1": main_topic[0].upper() if main_topic else "Type 1",
                "timeline": "2019-2032",
                "cur": f"{cur.upper()} {value_in.upper()}",
            },
            29: {
                "2_heading": headline_2,
                "timeline": f"HISTORIC YEAR {historical_year} FORECAST TO {forecast_year}",
            },
            30: {"heading": headline},
            31: {
                "company": company_info["company_name"].upper(),
                "e": company_info["employee_count"],
                "h": company_info["headquarters"],
                "geo": company_info["geographic_presence"],
                "es": company_info["estd"],
                "rev":company_info["revenue"],
            },
            32: {
                "2_heading": headline_2,
                "timeline": f"HISTORIC YEAR {historical_year} FORECAST TO {forecast_year}",
            },
            33: {
                "company": company_info["company_name"].upper(),
                "e": company_info["employee_count"],
                "ownership": company_info["ownership"],
                "h": company_info["headquarters"],
                "website": company_info["website"],
                "es": company_info["estd"],
                "product": company_info["top_product"],
                "para": company_info["short_description_company"],
                "rev":company_info["revenue"],
                "geo": company_info["geographic_presence"],
                "description": company_info["description_product"],
                },
            34: {"company": company_info["company_name"].upper()},
            35: {"company": company_info["company_name"].upper()},
        }

        logger.info(f"[{request_id}] Starting presentation modification")
        ppt_start_time = time.time()
        
        template_path = "testppt.pptx"
        if not os.path.exists(template_path):
            logger.error(f"[{request_id}] Template file not found: {template_path}")
            return jsonify({
                'error': 'Template file not found',
                'message': 'Please ensure testppt.pptx is in the project directory'
            }), 500
            
        logger.info(f"[{request_id}] Loading base presentation")
        prs = Presentation(template_path)
        
        for slide in prs.slides:
            data = slide_data.get(prs.slides.index(slide), {})
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for key, value in data.items():
                            token = f"{{{{{key}}}}}"
                            replace_text_in_paragraph(paragraph, token, value)

        for slide in prs.slides:
            data = slide_data.get(prs.slides.index(slide), {})
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for key, value in data.items():
                            token = f"{{{{{key}}}}}"
                            replace_text_preserving_color(paragraph, token, value)

        logger.info(f"[{request_id}] Adding taxonomy boxes")
        generator = TaxonomyBoxGenerator(prs)
        generator.add_taxonomy_boxes(1, table_taxonomy)

        logger.info(f"[{request_id}] Performing text replacements in tables")
        table_slide_indices = [10,13, 16, 17,18,19,21, 23, 24, 25, 26,27, 28,29,30,31, 32, 33,34,35]
        replace_text_in_tables(prs, table_slide_indices, slide_data)

        logger.info(f"[{request_id}] Creating Table of Contents")
        toc_slide_indices = [4, 5, 6, 7]
        add_toc_to_slides(prs, toc_data_levels, toc_slide_indices)

        logger.info(f"[{request_id}] Adding tables and charts")
        target_slide_indices = [24, 27, 28]
        graph_table = list(nested_dict[main_topic[0]].keys()) if main_topic else []
        total_rows = len(graph_table)
        
        row_labels = graph_table.copy() 
        row_labels.append("Total")

        years = [str(y) for y in range(2019, 2033)]
        columns = [""] + years + ["CAGR (2025–2032)"]
        num_rows = len(row_labels) + 1
        num_cols = len(columns)

        header_rgb = RGBColor(49, 6, 126)
        border_rgb = RGBColor(166, 166, 166)
        alt_row_colors = [RGBColor(231, 231, 231), RGBColor(255, 255, 255)]

        font_mapping = {
            "header": "Poppins Bold",
            "first_col": "Poppins Bold",
            "values": "Poppins Medium",
        }

        for slide_index in target_slide_indices:
            if slide_index < len(prs.slides):
                slide = prs.slides[slide_index]

                left = Inches(0.45)
                top = Inches(4.05)
                width = Inches(8.7)
                height = Inches(0.72 + num_rows * 0.3)
                table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

                for col_index, header in enumerate(columns):
                    cell = table.cell(0, col_index)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = header_rgb

                    cell.text_frame.clear()
                    para = cell.text_frame.paragraphs[0]
                    para.text = header.replace("\n", " ").strip()
                    
                    para.alignment = PP_ALIGN.CENTER
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE  
                    
                    if para.runs:
                        run = para.runs[0]
                    else:
                        run = para.add_run()
                    
                    if col_index != num_cols - 1:
                        run.font.size = Pt(5.7)
                        cell.text_frame.word_wrap = False
                    else:
                        run.font.size = Pt(8)
                    
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    run.font.name = font_mapping["header"]


                for row_index, label in enumerate(row_labels, start=1):
                    row_color = alt_row_colors[(row_index - 1) % 2]

                    for col_index in range(num_cols):
                        cell = table.cell(row_index, col_index)

                        if col_index == 0:
                            cell.text = label
                        elif col_index == num_cols - 1:
                            cell.text = "XX%"
                        else:
                            cell.text = "XX"

                        para = cell.text_frame.paragraphs[0]
                        para.alignment = PP_ALIGN.CENTER
                        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

                        if col_index == 0:
                            para.font.size = Pt(8)
                            para.font.name = font_mapping["first_col"]
                            para.font.bold = True
                        else:
                            para.font.size = Pt(9)
                            para.font.name = font_mapping["values"]

                            if label == "Total" and col_index == num_cols - 1:
                                para.font.bold = True
                            if row_index == num_rows - 1:
                                para.font.bold = True

                        cell.fill.solid()
                        cell.fill.fore_color.rgb = row_color
                        set_cell_border(cell)

                for col_index in range(num_cols):
                    if col_index == 0:
                        table.columns[col_index].width = Inches(1)
                    elif col_index == num_cols - 1:
                        table.columns[col_index].width = Inches(0.8)
                    else:
                        table.columns[col_index].width = Inches(0.4)

        if main_topic:
            chart_columns = graph_table

            for idx in target_slide_indices:
                if idx < len(prs.slides):
                    slide = prs.slides[idx]
                    data = generate_actual_data()
                    
                    create_chart_on_slide(
                        slide, data, chart_columns,
                        Inches(0.4), Inches(1.1), Inches(12.5), Inches(2.8)
                    )

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        clean_market_name = clean_filename(headline)
        filename = f"{clean_market_name}_{timestamp}.pptx"
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        
        logger.info(f"[{request_id}] Saving presentation to: {filepath}")
        prs.save(filepath)
        
        ppt_elapsed = time.time() - ppt_start_time
        total_elapsed = time.time() - start_time
        
        timing_logger.info(f"[{request_id}] PPT processing completed in {ppt_elapsed:.2f}s")
        timing_logger.info(f"[{request_id}] Total request completed in {total_elapsed:.2f}s")
        
        logger.info(f"[{request_id}] PPT generation completed successfully: {filename}")
        
        return jsonify({
            'success': True,
            'filename': filename,
            'message': 'PowerPoint generated successfully'
        })
        
    except Exception as e:
        elapsed = time.time() - start_time
        logger.error(f"[{request_id}] PPT generation failed after {elapsed:.2f}s: {str(e)}")
        traceback.print_exc()
        return jsonify({
            'error': 'Failed to generate PowerPoint',
            'message': str(e)
        }), 500


@app.route('/download')
def download_file(filename):
    download_start_time = time.time()
    logger.info(f"Download request for file: {filename}")
    
    try:
        import urllib.parse
        decoded_filename = urllib.parse.unquote(filename)
        
        if '..' in decoded_filename or '/' in decoded_filename or '\\' in decoded_filename:
            logger.warning(f"Invalid filename attempted: {decoded_filename}")
            return jsonify({'error': 'Invalid filename'}), 400
        
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], decoded_filename)
        
        if not os.path.exists(filepath):
            logger.warning(f"File not found: {filepath}")
            try:
                available_files = os.listdir(app.config['UPLOAD_FOLDER'])
                logger.info(f"Available files: {available_files}")
            except:
                logger.error("Could not list available files")
            return jsonify({'error': 'File not found'}), 404
        
        def remove_file_after_send(filepath):
            def remove_file(response):
                try:
                    if os.path.exists(filepath):
                        os.remove(filepath)
                        logger.info(f"Temporary file deleted: {filepath}")
                except Exception as e:
                    logger.error(f"Error deleting file: {e}")
                return response
            return remove_file
        
        response = send_file(filepath, as_attachment=True, download_name=decoded_filename)
        
        @response.call_on_close
        def delete_file():
            try:
                if os.path.exists(filepath):
                    os.remove(filepath)
                    logger.info(f"Temporary file deleted after download: {filepath}")
            except Exception as e:
                logger.error(f"Error deleting file after download: {e}")
        
        elapsed = time.time() - download_start_time
        timing_logger.info(f"File download completed in {elapsed:.2f}s: {decoded_filename}")
        
        return response
        
    except Exception as e:
        elapsed = time.time() - download_start_time
        logger.error(f"Download failed after {elapsed:.2f}s: {e}")
        return jsonify({'error': 'File download failed', 'details': str(e)}), 500


if __name__ == '__main__':
    if not os.path.exists('.env'):
        with open('.env', 'w') as f:
            f.write('OPENAI_API_KEY=your_openai_api_key_here\n')
            f.write('GEMINI_API_KEY=your_gemini_api_key_here\n')
        logger.info("Created .env file. Please add your API keys.")
    
    if not os.path.exists('templates/index.html'):
        os.makedirs('templates', exist_ok=True)
        logger.warning("Please save the HTML content from the artifact to templates/index.html")
    
    logger.info("Starting Flask application")
    app.run(host="0.0.0.0",debug=True, port=5000,)